from typing import Optional, Literal, List

import pptx
import pptx.slide, pptx.dml, pptx.dml.color, pptx.util, pptx.oxml

from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


SlideType = pptx.slide.Slide


class Slide(object):

    def __init__(self, slide: SlideType, title: Optional[str] = None):
        self.slide: SlideType = slide

        self.title: Optional[str] = title
        if self.title is not None:
            self.writeTitle()

    def setTitle(self, title: str) -> None:
        self.title = title

    def writeTitle(self) -> None:

        self.addTextbox(self.title,
                        height=0.44, width=4.06,
                        left=0.35, top=0.25,
                        font='Calibri', fontsize=20,
                        bold=True, color='black',
                        alignment='left', border=False)

    def addImage(self, picpath: str, height=None, width=None, left=None, top=None, border=False, raw_dims=False):

        if not raw_dims:
            height, width, left, top = Slide.convertToLocalInches(height, width, left, top)

        picture = self.slide.shapes.add_picture(
            picpath,
            height=height,
            width=width,
            left=left,
            top=top,
        )

        if border:
            line = picture.line
            line.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
            line.width = pptx.util.Inches(0.01)

    def addTextbox(self, text, height=None, width=None, left=None, top=None,
                   border=False, font=None, fontsize=20, bold=False, italic=False, color='black', alignment='center',
                   raw_dims=False):

        if not raw_dims:
            height, width, left, top = Slide.convertToLocalInches(height, width, left, top)

        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.paragraphs[0].text = text

        if not raw_dims:
            fontsize = pptx.util.Pt(fontsize)

        tf.paragraphs[0].font.size = fontsize
        tf.paragraphs[0].font.name = font

        # Si viene como 'red', 'blue', etc.
        if isinstance(color, str):
            # raise ValueError("Implementar clase de colores")
            # color = [int(v * 255) for v in mm_colors.colors_dict[color]['rgb']]
            color = [0, 0, 0]
            tf.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(*color)
        else:
            tf.paragraphs[0].font.color.rgb = color

        if border:
            txBox.line.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)

        if bold:
            tf.paragraphs[0].font.bold = True

        if italic:
            tf.paragraphs[0].font.italic = True

        alignment_dict = {
            'center': PP_ALIGN.CENTER,
            'left': PP_ALIGN.LEFT,
            'right': PP_ALIGN.RIGHT,
        }

        try:
            tf.paragraphs[0].alignment = alignment_dict[alignment]
        except:
            pass

    def addTable(self, data_dict, keys, height=None, width=None, left=None, top=None,
                 orientation: Literal['rows', 'columns'] = 'rows', fontsize=10,
                 col_width=None, col_widths_list=None):

        height, width, left, top = Slide.convertToLocalInches(height, width, left, top)

        if orientation == 'rows':
            shape = self.slide.shapes.add_table(len(keys), 2, left, top, width, height)
        else:
            test_key = list(data_dict.keys())[0]
            shape = self.slide.shapes.add_table(len(data_dict[test_key]) + 1, len(keys), left, top, width, height)

        Slide.setTableWidths(shape.table, data_dict, keys, orientation=orientation,
                             col_width=col_width, col_widths_list=col_widths_list)
        Slide.setTableFormat(shape.table, data_dict, keys, orientation=orientation)

        if orientation == 'rows':

            for i, k in enumerate(keys):
                cell = shape.table.cell(i, 0)
                cell.text = k
                cell.text_frame.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
                cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(fontsize)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()

                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                cell = shape.table.cell(i, 1)
                cell.text = str(data_dict[k])
                cell.text_frame.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
                cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(fontsize)
                cell.fill.solid()

                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        elif orientation == 'columns':

            for i, k in enumerate(keys):
                cell = shape.table.cell(0, i)
                cell.text = str(k)
                cell.text_frame.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
                cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(fontsize)
                cell.fill.solid()
                cell.text_frame.paragraphs[0].font.bold = True

                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                cell.margin_left = 0
                cell.margin_right = 0
                cell.margin_top = 0
                cell.margin_bottom = 0

                for j, v in enumerate(data_dict[k]):
                    cell = shape.table.cell(j + 1, i)
                    cell.text = str(v)
                    cell.text_frame.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
                    cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(fontsize)
                    cell.fill.solid()

                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                    cell.margin_left = 0
                    cell.margin_right = 0
                    cell.margin_top = 0
                    cell.margin_bottom = 0

    def addHyperlink(self, text, slide_destination, height=None, width=None, left=None, top=None,
                     font='Calibri', fontsize=20, bold=False, italic=False, alignment='center'):

        height, width, left, top = Slide.convertToLocalInches(height, width, left, top)

        txBox = self.slide.shapes.add_textbox(left, top, width, height)

        p = txBox.text_frame.paragraphs[0]

        p.font.size = pptx.util.Pt(fontsize)
        p.font.name = font

        if bold:
            p.font.bold = True

        if italic:
            p.font.italic = True

        alignment_dict = {
            'center': PP_ALIGN.CENTER,
            'left': PP_ALIGN.LEFT,
            'right': PP_ALIGN.RIGHT,
        }

        p.alignment = alignment_dict[alignment]

        r = p.add_run()
        r.text = text

        # Hyperlink
        rId = self.slide.part.relate_to(slide_destination.part, RT.SLIDE)
        rPr = r._r.get_or_add_rPr()
        hlinkClick = rPr.add_hlinkClick(rId)
        hlinkClick.set('action', 'ppaction://hlinksldjump')

    # Move all the below methods TODO
    @staticmethod
    def convertToLocalInches(height, width, left, top):
        height = pptx.util.Inches(height) if height is not None else None
        width = pptx.util.Inches(width) if width is not None else None
        left = pptx.util.Inches(left) if left is not None else None  # Estos no deberian ser nunca None
        top = pptx.util.Inches(top) if top is not None else None  # Estos no deberian ser nunca None

        if top is None or left is None:
            raise ValueError('Both top and left should be specified.')

        return height, width, left, top

    @staticmethod
    def SubElement(parent, tagname, **kwargs):
        element = pptx.oxml.xmlchemy.OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    @staticmethod
    def _set_cell_border(cell, border_color="000000", border_width='12700'):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        for lines in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            ln = Slide.SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
            solidFill = Slide.SubElement(ln, 'a:solidFill')
            srgbClr = Slide.SubElement(solidFill, 'a:srgbClr', val=border_color)
            prstDash = Slide.SubElement(ln, 'a:prstDash', val='solid')
            round_ = Slide.SubElement(ln, 'a:round')
            headEnd = Slide.SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
            tailEnd = Slide.SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

    @staticmethod
    def setTableWidths(table, data_dict, keys, orientation='rows', col_width=None,
                       col_widths_list: Optional[List[float]] = None):

        if orientation == 'rows':
            if col_widths_list is not None:
                for i, w in enumerate(col_widths_list):
                    table.columns[i].width = pptx.util.Inches(w)
            else:
                table.columns[0].width = pptx.util.Inches(1.)
                table.columns[1].width = pptx.util.Inches(0.6)

        elif orientation == 'columns':
            for i in range(len(keys)):

                if col_width is not None:
                    table.columns[i].width = pptx.util.Inches(col_width)
                else:
                    table.columns[i].width = pptx.util.Inches(col_widths_list[i])

    @staticmethod
    def setTableFormat(table, data_dict, keys, orientation='rows'):

        for i, k in enumerate(keys):

            if orientation == 'rows':
                cell = table.cell(i, 0)
                Slide._set_cell_border(cell, border_color="000000", border_width='12700')

                cell.fill.solid()
                cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(255, 255, 255)

                cell = table.cell(i, 1)
                Slide._set_cell_border(cell, border_color="000000", border_width='12700')

                cell.fill.solid()
                cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(255, 255, 255)


            else:
                for j in range(len(data_dict[k]) + 1):
                    cell = table.cell(j, i)
                    Slide._set_cell_border(cell, border_color="000000", border_width='12700')

                    cell.fill.solid()
                    cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(255, 255, 255)
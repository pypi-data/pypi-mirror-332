import os
import sys
import datetime
from typing import Optional, List, Literal

import pptx
import pptx.slide
import pptx.presentation

from pptx.util import Inches

# import mymesh.mymesh.variables.colors as mm_colors

from geoassistant.report.Slide import Slide
from geoassistant.report.SlidesFileTemplates import SlidesFileTemplates

PresentationType = pptx.presentation.Presentation


class SlidesFile(object):

    def __init__(self, style: Literal['default', 'GMT', 'GMT_vertical'] = 'default'):

        self.presentation: Optional[PresentationType] = None

        self.master_slides = {}

        self.slides: List['Slide'] = []
        self.front_page: Optional['Slide'] = None

        if style == 'default':
            self.initiateEmptyPresentation()
        else:
            self.loadStyle(style=style)

        self.templates: SlidesFileTemplates = SlidesFileTemplates(self)

    def setPresentation(self, presentation: PresentationType) -> None:
        self.presentation = presentation

    def getPresentation(self) -> pptx.presentation.Presentation:
        return self.presentation

    def getSlides(self) -> List['Slide']:
        return self.slides

    @staticmethod
    def readPPTFile(filepath: str) -> 'SlidesFile':
        sf = SlidesFile()
        sf.setPresentation(pptx.Presentation(filepath))
        sf._processExistingSlides()
        return sf

    def initiateEmptyPresentation(self) -> None:
        self.presentation = pptx.Presentation()
        self.presentation.slide_width = Inches(13.333)
        self.presentation.slide_height = Inches(7.5)

    def loadStyle(self, style: Literal['GMT', 'GMT_vertical']) -> None:
        raise ValueError("IMPLEMENTAR")
        # base_path = os.environ.get('PYTHONPATH') + '/reports/bases/'
        base_path = "C:/Users/jorge/OneDrive/Desktop/JMV/04. Librerias/Python/" + '/reports/bases/'
        base_path += 'base_' + style + '.pptx'

        self.presentation = pptx.Presentation(base_path)
        self._processExistingSlides()

    def _processExistingSlides(self) -> None:
        for s in self.presentation.slides:
            self.slides += [Slide(s)]

    def addEmptySlide(self, title: str = '{TITULO}'):

        blank_slide_layout = self.presentation.slide_layouts[1]
        s = self.presentation.slides.add_slide(blank_slide_layout)
        # Saco el cuerpo predeterminado
        for i in [0, 1]:
            sp = s.shapes[0].element
            sp.getparent().remove(sp)

        slide = Slide(s, title=title)
        self.slides += [slide]

        return slide

    def savePresentation(self, savepath: str) -> None:
        self.presentation.save(savepath)

    def create(self, images_paths: List[str],
               images_grid: List[int],
               images_titles: Optional[List[str]] = None,
               images_comments: Optional[List[str]] = None):
        self.templates.create(images_paths, images_grid, images_titles=images_titles, images_comments=images_comments)

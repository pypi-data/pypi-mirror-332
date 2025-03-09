import json, urllib, cv2, os, shutil
from pdf2image import convert_from_path
import pytesseract
import PyPDF2
from docx import Document
from pptx import Presentation
# import image_utils, file_utils, text_utils
from pii_scanner.octopii import image_utils, file_utils, text_utils

model_file_name = 'models/other_pii_model.h5'
labels_file_name = 'models/other_pii_model.txt'
temp_dir = ".OCTOPII_TEMP/"
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '3'

def extract_text(file_path):
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    if file_extension in ['.pdf']:
        text = ""
        with open(file_path, 'rb') as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            for page in reader.pages:
                text += page.extract_text()

    elif file_extension in ['.docx']:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])

    elif file_extension in ['.pptx']:
        presentation = Presentation(file_path)
        text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])

    elif file_extension in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
        image = cv2.imread(file_path)
        text = pytesseract.image_to_string(image)

    elif file_extension in ['.txt']:
        with open(file_path, 'r', encoding='utf-8') as txt_file:
            text = txt_file.read()

    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

    return text

def search_pii(file_path):
    contains_faces = 0
    if file_utils.is_image(file_path):
        image = cv2.imread(file_path)
        contains_faces = image_utils.scan_image_for_people(image)
        original, intelligible = image_utils.scan_image_for_text(image)
        text = original

    elif file_utils.is_pdf(file_path):
        pdf_pages = convert_from_path(file_path, 400)  # Higher DPI reads small text better
        for page in pdf_pages:
            contains_faces = image_utils.scan_image_for_people(page)
            original, intelligible = image_utils.scan_image_for_text(page)
            text = original

    else:
        text = extract_text(file_path)
        intelligible = text_utils.string_tokenizer(text)
        
    rules = text_utils.get_regexes()
    
    addresses = text_utils.regional_pii(text)
    emails = text_utils.email_pii(text, rules)
    phone_numbers = text_utils.phone_pii(text, rules)

    keywords_scores = text_utils.keywords_classify_pii(rules, intelligible)
    score = max(keywords_scores.values())
    pii_class = list(keywords_scores.keys())[list(keywords_scores.values()).index(score)]

    country_of_origin = rules[pii_class]["region"]

    identifiers = text_utils.id_card_numbers_pii(text, rules)

    if score < 5:
        pii_class = None

    if len(identifiers) != 0:
        identifiers = identifiers[0]["result"]

    if temp_dir in file_path:
        file_path = file_path.replace(temp_dir, "")
        file_path = urllib.parse.unquote(file_path)

    result = {
        "file_path": file_path,
        "pii_class": pii_class,
        "score": score,
        "country_of_origin": country_of_origin,
        "faces": contains_faces,
        "identifiers": identifiers,
        "emails": emails,
        "phone_numbers": phone_numbers,
        "addresses": addresses
    }

    return result

def process_file_octopii(location, output_file="output.json", notifyURL=None):
    rules = text_utils.get_regexes()
    files = []
    temp_exists = False

    try:
        shutil.rmtree(temp_dir)
    except:
        pass

    if "http" in location:
        try:
            file_urls = []
            _, extension = os.path.splitext(location)
            if extension != "":
                file_urls.append(location)
            else:
                files = file_utils.list_local_files(location)

            file_urls = file_utils.list_s3_files(location)
            if len(file_urls) != 0:
                temp_exists = True
                os.makedirs(os.path.dirname(temp_dir))
                for url in file_urls:
                    file_name = urllib.parse.quote(url, "UTF-8")
                    urllib.request.urlretrieve(url, temp_dir + file_name)

        except:
            try:
                file_urls = file_utils.list_directory_files(location)

                if len(file_urls) != 0:  # directory listing (e.g.: Apache)
                    temp_exists = True
                    os.makedirs(os.path.dirname(temp_dir))
                    for url in file_urls:
                        try:
                            encoded_url = urllib.parse.quote(url, "UTF-8")
                            urllib.request.urlretrieve(url, temp_dir + encoded_url)
                        except:
                            pass  # capture 404

                else:  # curl text from location if available
                    temp_exists = True
                    os.makedirs(os.path.dirname(temp_dir))
                    encoded_url = urllib.parse.quote(location, "UTF-8") + ".txt"
                    urllib.request.urlretrieve(location, temp_dir + encoded_url)

            except:
                traceback.print_exc()
                print("This URL is not a valid S3 or has no directory listing enabled. Try running Octopii on these files locally.")
                return None

        files = file_utils.list_local_files(temp_dir)

    else:
        _, extension = os.path.splitext(location)
        if extension != "":
            files.append(location)
        else:
            files = file_utils.list_local_files(location)

    if len(files) == 0:
        print("Invalid path provided. Please provide a non-empty directory or a file as an argument.")
        return None

    results_list = []

    # try and truncate files if they're too big
    for file_path in files:
        try:
            file_utils.truncate(file_path)
        except:
            pass

    for file_path in files:
        try:
            results = search_pii(file_path)
            results_list.append(results)
            file_utils.append_to_output_file(results, output_file)
            if notifyURL is not None:
                webhook.push_data(json.dumps(results), notifyURL)

        except ValueError as e:
            print(f"Error processing file '{file_path}': {str(e)}")

    if temp_exists:
        shutil.rmtree(temp_dir)

    return results_list


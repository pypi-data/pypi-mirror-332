from typing import List

def extract_text_from_slides(pptx_path: str) -> List[List[List[str]]]:
    """Extract text from a PowerPoint presentation file and return it as a list of lists of strings."""
    from pptx import Presentation

    prs = Presentation(pptx_path)
    all_texts = []

    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                shape_data = []
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        shape_data.append(run.text)
                slide_texts.append(shape_data)
            elif hasattr(shape, "text"):
                slide_texts.append(shape.text)
        all_texts.append(slide_texts)

    return all_texts


def replace_text_in_slides(pptx_path: str, new_texts: List[List[List[str]]], output_path: str):
    """Replace text in a PowerPoint presentation file with new text."""
    from pptx import Presentation

    prs = Presentation(pptx_path)

    for slide, slide_texts in zip(prs.slides, new_texts):
        shape_index = 0
        for shape in slide.shapes:
            print("Shape", shape_index)
            if hasattr(shape, "text_frame"):
                print("A")
                text_data = slide_texts[shape_index]
                print("text_data", text_data)

                run_index = 0
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = translate(run.text, text_data[run_index])
                        run_index += 1
                shape_index += 1
            elif hasattr(shape, "text"):
                print("B")
                shape.text = translate(shape.text, slide_texts[shape_index])
                shape_index += 1

    prs.save(output_path)

def translate(old_text, new_text):
    if old_text != new_text:
        print(f"Translating '{old_text}' to '{new_text}'")
    return new_text